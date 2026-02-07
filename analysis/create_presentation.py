"""
Create PowerPoint Presentation - RM Results
Plots take maximum space, aspect ratios preserved
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# Slide dimensions (inches) - Standard 16:9
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# Margins
MARGIN = 0.2

# Slide content
SLIDES = [
    {
        "title": "Exp1: Number Deviation",
        "image": "outputs/exp1/figures/number_deviation.png",
        "bullets": [
            "Reported vs actual set size (3, 4, 5 bars)",
            "Black dot = grand mean with 95% CI",
            "Result: Strong RM effect - participants consistently underreport"
        ]
    },
    {
        "title": "Exp1: Spacing Deviation",
        "image": "outputs/exp1/figures/spacing_deviation.png",
        "bullets": [
            "By RM condition, faceted by spacing category",
            "Teal = Non-RM; Magenta = RM",
            "Result: RM trials show larger (overestimated) spacing"
        ]
    },
    {
        "title": "Exp1: Width Deviation",
        "image": "outputs/exp1/figures/width_deviation.png",
        "bullets": [
            "By RM condition, faceted by bar width",
            "Stars = significant differences (p < .05)",
            "Result: RM trials show wider, more accurate width perception"
        ]
    },
    {
        "title": "Exp1: Bayesian Width",
        "image": "outputs/exp1/figures/bayesian_width_deviation.png",
        "bullets": [
            "Bayesian one-sample t-test: deviation from 0?",
            "Grey = Prior; Colored = Posterior",
            "Result: Non-RM strongly deviates from 0; RM closer to accurate"
        ]
    },
    {
        "title": "Exp1: Density Split",
        "image": "outputs/exp1/figures/density_category_split.png",
        "bullets": [
            "Density = (set size × width) / array extent",
            "Split into 3 equal groups: Low / Medium / High",
            "Defines grouping for density-based analysis"
        ]
    },
    {
        "title": "Exp1: Density Deviation",
        "image": "outputs/exp1/figures/density_deviation.png",
        "bullets": [
            "By RM condition across density categories",
            "Error bars = 95% CI from participant means",
            "Result: Density perception similar in RM and Non-RM trials"
        ]
    },
    {
        "title": "Exp1: Density Bayesian",
        "image": "outputs/exp1/figures/density_deviation_ridgeplot.png",
        "bullets": [
            "Posterior distributions for RM vs Non-RM difference",
            "BF01 values indicate evidence for no difference",
            "Result: No reliable density perception difference between conditions"
        ]
    },
    {
        "title": "Exp2: Number Deviation",
        "image": "outputs/exp2/figures/number_deviation.png",
        "bullets": [
            "Exp2 uses single-bar probe (different from Exp1)",
            "Same pattern of underreporting as Exp1",
            "Result: RM effect replicates with different probe method"
        ]
    },
    {
        "title": "Exp2: Width Deviation",
        "image": "outputs/exp2/figures/width_deviation.png",
        "bullets": [
            "Single-bar probe: no multi-bar context",
            "Smaller effect compared to Exp1",
            "Result: RM width advantage reduced without multi-bar probe"
        ]
    },
    {
        "title": "Exp2: Bayesian Width",
        "image": "outputs/exp2/figures/bayesian_width_deviation.png",
        "bullets": [
            "Same Bayesian analysis as Exp1 for comparison",
            "BF01 values shown for each condition",
            "Result: Both conditions deviate; RM still closer to accurate"
        ]
    },
    {
        "title": "Combined Comparison",
        "image": "outputs/combined_width_deviation_barplot.png",
        "bullets": [
            "Exp1 (multi-bar probe) vs Exp2 (single-bar probe)",
            "Relative width deviation by experiment and RM condition",
            "Result: Multi-bar probe amplifies the RM width advantage"
        ]
    }
]


def get_image_size(path):
    """Get image dimensions and aspect ratio"""
    with Image.open(path) as img:
        width, height = img.size
        return width, height, width / height


def calculate_image_placement(aspect_ratio, is_wide_layout):
    """Calculate image size and position preserving aspect ratio"""

    if is_wide_layout:
        # Wide image: use most of slide, leave space at bottom for bullets
        max_width = SLIDE_WIDTH - 2 * MARGIN
        max_height = SLIDE_HEIGHT - 1.8  # Leave ~1.8 inches for bullets at bottom

        # Calculate size preserving aspect ratio
        if max_width / aspect_ratio <= max_height:
            img_width = max_width
            img_height = max_width / aspect_ratio
        else:
            img_height = max_height
            img_width = max_height * aspect_ratio

        # Center horizontally, top aligned
        img_left = (SLIDE_WIDTH - img_width) / 2
        img_top = MARGIN

        # Bullets below image
        bullet_left = MARGIN
        bullet_top = img_top + img_height + 0.15
        bullet_width = SLIDE_WIDTH - 2 * MARGIN

    else:
        # Square/tall image: image on left, bullets on right
        max_width = SLIDE_WIDTH * 0.65
        max_height = SLIDE_HEIGHT - 2 * MARGIN

        # Calculate size preserving aspect ratio
        if max_width / aspect_ratio <= max_height:
            img_width = max_width
            img_height = max_width / aspect_ratio
        else:
            img_height = max_height
            img_width = max_height * aspect_ratio

        # Left aligned, vertically centered
        img_left = MARGIN
        img_top = (SLIDE_HEIGHT - img_height) / 2

        # Bullets on right
        bullet_left = img_left + img_width + 0.2
        bullet_top = SLIDE_HEIGHT * 0.3
        bullet_width = SLIDE_WIDTH - bullet_left - MARGIN

    return {
        "img_left": img_left,
        "img_top": img_top,
        "img_width": img_width,
        "img_height": img_height,
        "bullet_left": bullet_left,
        "bullet_top": bullet_top,
        "bullet_width": bullet_width
    }


def add_slide(prs, slide_data):
    """Add a slide with image and bullets"""

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    image_path = slide_data["image"]

    if not os.path.exists(image_path):
        print(f"Warning: {image_path} not found")
        return

    # Get image info
    _, _, aspect_ratio = get_image_size(image_path)
    is_wide = aspect_ratio > 1.4

    # Calculate placement
    placement = calculate_image_placement(aspect_ratio, is_wide)

    # Add image
    slide.shapes.add_picture(
        image_path,
        Inches(placement["img_left"]),
        Inches(placement["img_top"]),
        Inches(placement["img_width"]),
        Inches(placement["img_height"])
    )

    # Add title (small, top-right)
    title_box = slide.shapes.add_textbox(
        Inches(SLIDE_WIDTH - 3.5),
        Inches(0.1),
        Inches(3.3),
        Inches(0.4)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data["title"]
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.RIGHT

    # Add bullets
    bullet_box = slide.shapes.add_textbox(
        Inches(placement["bullet_left"]),
        Inches(placement["bullet_top"]),
        Inches(placement["bullet_width"]),
        Inches(2)
    )
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True

    for i, bullet in enumerate(slide_data["bullets"]):
        if i == 0:
            para = bullet_frame.paragraphs[0]
        else:
            para = bullet_frame.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(14)
        para.space_after = Pt(4)


def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # Add slides
    for slide_data in SLIDES:
        add_slide(prs, slide_data)
        print(f"Added: {slide_data['title']}")

    # Save
    output_path = "outputs/rm_results_presentation.pptx"
    prs.save(output_path)
    print(f"\nSaved to: {output_path}")
    print(f"Total slides: {len(SLIDES)}")


if __name__ == "__main__":
    main()
